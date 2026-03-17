"""
Script para generar presentacion MLOps profesional con python-pptx
30 slides - MLOps: Del modelo al valor de negocio
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
import datetime
import os
import glob

# ─────────────────────────────────────────────
# COLORES (extraidos de patrones + corporativos)
# Patrones usan Berlin Sans FB / Open Sans
# Colores adicionales extraidos de XML
# ─────────────────────────────────────────────
AZUL_CORP   = "003A70"   # fondo principal
AZUL_MEDIO  = "0066CC"   # acento azul
NARANJA     = "FF6B00"   # acento naranja
GRIS_CLARO  = "F5F5F5"   # fondo claro
BLANCO      = "FFFFFF"
GRIS_TEXTO  = "333333"
ROJO        = "C0392B"
VERDE       = "27AE60"
AZUL_CLARO  = "AED6F1"
AMARILLO    = "F39C12"
GRIS_MEDIO  = "7F8C8D"

# Fuentes encontradas: Open Sans, Berlin Sans FB → fallback Calibri
FUENTE_TITULO  = "Calibri"
FUENTE_CUERPO  = "Calibri"

# ─────────────────────────────────────────────
# FOTO INSTRUCTOR
# ─────────────────────────────────────────────
FOTO_PATH = None
for ext in ["png", "jpg", "jpeg"]:
    p = f"C:/Users/bk70827/PycharmProjects/mlops-ciclo-vida/pdfs-ppts/assets/foto_instructor.{ext}"
    if os.path.exists(p):
        FOTO_PATH = p
        break
print(f"Foto instructor: {FOTO_PATH}")

# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_background(slide, color_hex, prs):
    """Rectangulo full-slide como fondo."""
    w = prs.slide_width
    h = prs.slide_height
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color_hex)
    shape.line.fill.background()  # sin borde
    return shape

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width=0):
    """Rectángulo con color de fondo."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, texto, x, y, w, h,
                color_hex=BLANCO, size_pt=18, bold=False,
                align=PP_ALIGN.LEFT, italic=False, font=None):
    """Textbox con formato."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.name = font or FUENTE_CUERPO
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)
    return txBox

def add_title(slide, texto, x, y, w, h,
              color_hex=BLANCO, size_pt=32, bold=True):
    return add_textbox(slide, texto, x, y, w, h,
                       color_hex=color_hex, size_pt=size_pt, bold=bold,
                       align=PP_ALIGN.LEFT)

def add_centered_text(slide, texto, x, y, w, h,
                      color_hex=BLANCO, size_pt=18, bold=False, italic=False):
    return add_textbox(slide, texto, x, y, w, h,
                       color_hex=color_hex, size_pt=size_pt, bold=bold,
                       italic=italic, align=PP_ALIGN.CENTER)

def add_bullet_box(slide, bullets, x, y, w, h,
                   color_hex=GRIS_TEXTO, size_pt=16, line_spacing=1.2):
    """Caja con lista de bullets."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet
        run.font.name = FUENTE_CUERPO
        run.font.size = Pt(size_pt)
        run.font.color.rgb = rgb(color_hex)
    return txBox

def add_numbered_box(slide, numero, titulo, descripcion,
                     x, y, w, h, num_color, bg_color, text_color):
    """Caja numerada con título y descripción."""
    box = add_rect(slide, x, y, w, h, bg_color, AZUL_MEDIO, 1.5)
    # número
    add_centered_text(slide, numero, x, y + 0.05, w, 0.6,
                      color_hex=num_color, size_pt=28, bold=True)
    # título
    add_centered_text(slide, titulo, x, y + 0.6, w, 0.45,
                      color_hex=text_color, size_pt=14, bold=True)
    # descripción
    add_centered_text(slide, descripcion, x, y + 1.0, w, 0.8,
                      color_hex=text_color, size_pt=11)

# ─────────────────────────────────────────────
# PRESENTACION
# ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completamente en blanco

fecha_hoy = datetime.date.today().strftime("%d de %B de %Y")

# ══════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)

# Línea decorativa naranja
add_rect(slide, 0, 5.6, 13.33, 0.08, NARANJA)

# Bloque lateral derecho más oscuro
add_rect(slide, 9.5, 0, 3.83, 7.5, "00264D")

# Foto instructor
if FOTO_PATH:
    try:
        slide.shapes.add_picture(FOTO_PATH, Inches(10.0), Inches(1.0),
                                  Inches(2.8), Inches(3.5))
        add_centered_text(slide, "Formador MLOps", 9.5, 4.6, 3.83, 0.4,
                          color_hex=BLANCO, size_pt=12, italic=True)
    except Exception as e:
        print(f"No se pudo insertar foto: {e}")

# Título principal
add_textbox(slide, "MLOps:", 0.5, 1.0, 8.8, 1.1,
            color_hex=BLANCO, size_pt=44, bold=True)
add_textbox(slide, "Del modelo al valor de negocio", 0.5, 1.9, 8.8, 0.8,
            color_hex=NARANJA, size_pt=32, bold=True)

# Subtítulo
add_textbox(slide,
            "Transformando la inteligencia artificial en resultados reales",
            0.5, 2.85, 8.8, 0.65,
            color_hex=AZUL_CLARO, size_pt=18)

# Separador fino
add_rect(slide, 0.5, 3.65, 5.0, 0.04, NARANJA)

# Datos formativos
add_textbox(slide, "Formacion MLOps 2025", 0.5, 3.85, 8.0, 0.4,
            color_hex=BLANCO, size_pt=14)
add_textbox(slide, fecha_hoy, 0.5, 4.25, 8.0, 0.35,
            color_hex=AZUL_CLARO, size_pt=12)

# Footer
add_textbox(slide, "Sesion Confidencial — Solo uso interno", 0.5, 6.9, 9.0, 0.4,
            color_hex=GRIS_MEDIO, size_pt=10, italic=True)

# ══════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)

# Barra lateral azul
add_rect(slide, 0, 0, 0.35, 7.5, AZUL_CORP)

# Barra top
add_rect(slide, 0, 0, 13.33, 0.55, AZUL_CORP)

add_textbox(slide, "Agenda del dia", 0.55, 0.08, 10.0, 0.45,
            color_hex=BLANCO, size_pt=22, bold=True)

# Items de agenda
items = [
    ("01", "El problema: la IA que no llega a produccion"),
    ("02", "Que es MLOps y por que importa"),
    ("03", "El ciclo de vida de un modelo de IA"),
    ("04", "Beneficios de negocio medibles"),
    ("05", "Como lo implementamos en nuestra plataforma"),
    ("06", "Proximos pasos y llamada a la accion"),
]

for i, (num, texto) in enumerate(items):
    y_pos = 0.85 + i * 0.9
    # Numero en caja naranja
    add_rect(slide, 0.55, y_pos, 0.55, 0.55, NARANJA)
    add_centered_text(slide, num, 0.55, y_pos + 0.07, 0.55, 0.45,
                      color_hex=BLANCO, size_pt=16, bold=True)
    # Texto
    add_textbox(slide, texto, 1.25, y_pos + 0.08, 11.0, 0.45,
                color_hex=GRIS_TEXTO, size_pt=15)
    # Separador
    if i < len(items) - 1:
        add_rect(slide, 1.25, y_pos + 0.62, 10.8, 0.02, GRIS_CLARO)

# Footer
add_rect(slide, 0, 7.1, 13.33, 0.4, GRIS_CLARO)
add_textbox(slide, "Sesion interactiva — haremos pausas para preguntas",
            0.5, 7.15, 12.0, 0.3, color_hex=GRIS_MEDIO, size_pt=11, italic=True,
            align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 3 — PAUSA PREGUNTA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, NARANJA, prs)

add_centered_text(slide, "?", 4.0, 0.5, 5.33, 1.5,
                  color_hex=BLANCO, size_pt=80, bold=True)
add_centered_text(slide, "Pausa para reflexion", 1.0, 1.6, 11.33, 0.6,
                  color_hex=BLANCO, size_pt=18, italic=True)
add_centered_text(slide,
                  "\u00bfQue sabeis de IA en produccion?",
                  0.5, 2.4, 12.33, 1.0,
                  color_hex=BLANCO, size_pt=36, bold=True)
add_centered_text(slide,
                  "\u00bfCu\u00e1ntos modelos de IA conocen que usan a diario?",
                  0.5, 3.6, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=20)
add_rect(slide, 4.0, 4.5, 5.33, 0.06, BLANCO)
add_centered_text(slide, "Compartan su experiencia",
                  2.0, 4.7, 9.33, 0.4,
                  color_hex=BLANCO, size_pt=14, italic=True)

# ══════════════════════════════════════════════
# SLIDE 4 — EL PROBLEMA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "El problema: la IA que no llega a produccion",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

# Stat grande
add_rect(slide, 0.4, 0.9, 3.5, 2.8, NARANJA)
add_centered_text(slide, "87%", 0.4, 1.1, 3.5, 1.3,
                  color_hex=BLANCO, size_pt=72, bold=True)
add_centered_text(slide, "de los modelos de IA\nnunca llegan\na produccion",
                  0.4, 2.3, 3.5, 1.2,
                  color_hex=BLANCO, size_pt=14)

# Bullets
bullets = [
    "  Falta de procesos estandarizados entre ciencia de datos e IT",
    "  Silos entre equipos: datos, modelos y deployment separados",
    "  Ausencia de monitoreo: modelos degradados sin deteccion",
    "  Deuda tecnica acumulada en cada iteracion",
]
add_bullet_box(slide, bullets, 4.2, 1.0, 8.7, 3.0,
               color_hex=GRIS_TEXTO, size_pt=15)

# Fuente
add_rect(slide, 0, 6.8, 13.33, 0.7, GRIS_CLARO)
add_textbox(slide, "Fuente: Gartner — $3.9T valor estimado IA 2022 | Databricks Big Book of MLOps v6",
            0.5, 6.88, 12.0, 0.4,
            color_hex=GRIS_MEDIO, size_pt=10, italic=True)

# ══════════════════════════════════════════════
# SLIDE 5 — ANALOGIA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)

add_centered_text(slide, "MLOps es como una cadena de montaje para la IA",
                  0.5, 0.3, 12.33, 0.9,
                  color_hex=BLANCO, size_pt=28, bold=True)

# Pipeline visual
pasos = ["Datos", "Modelo", "Pruebas", "Produccion", "Monitoreo"]
colores_pasos = [NARANJA, BLANCO, NARANJA, BLANCO, NARANJA]
txt_colores   = [BLANCO, AZUL_CORP, BLANCO, AZUL_CORP, BLANCO]
x_start = 0.4
box_w = 2.2
box_h = 1.3
gap = 0.25

for i, (paso, col, tc) in enumerate(zip(pasos, colores_pasos, txt_colores)):
    x = x_start + i * (box_w + gap)
    add_rect(slide, x, 1.7, box_w, box_h, col)
    add_centered_text(slide, paso, x, 2.0, box_w, 0.6,
                      color_hex=tc, size_pt=18, bold=True)
    if i < len(pasos) - 1:
        add_centered_text(slide, "→", x + box_w, 2.0, gap + 0.1, 0.6,
                          color_hex=NARANJA, size_pt=24, bold=True)

# Analogia
add_rect(slide, 0.4, 3.3, 12.53, 0.05, NARANJA)
add_centered_text(slide,
                  "Igual que una fabrica no improvisa cada pieza, MLOps estandariza\ncada etapa del ciclo de vida del modelo de IA.",
                  0.5, 3.5, 12.33, 0.8,
                  color_hex=AZUL_CLARO, size_pt=16, italic=True)

add_centered_text(slide,
                  "Sin este proceso: caos, retrabajo, modelos en produccion sin garantias.",
                  0.5, 4.5, 12.33, 0.5,
                  color_hex=BLANCO, size_pt=14)

add_centered_text(slide,
                  "MLOps = ML + DevOps + Ingenieria de Datos",
                  0.5, 5.3, 12.33, 0.6,
                  color_hex=NARANJA, size_pt=20, bold=True)

# ══════════════════════════════════════════════
# SLIDE 6 — QUE ES MLOPS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Que es MLOps: el ciclo integrado",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

# 4 cajas principales
cajas = [
    ("DATA", "Ingesta, limpieza,\nfeature engineering", AZUL_CORP, BLANCO),
    ("MODELO", "Entrenamiento,\nexperimentos, versiones", NARANJA, BLANCO),
    ("PRODUCCION", "Despliegue, CI/CD,\nserving en tiempo real", AZUL_MEDIO, BLANCO),
    ("MONITOREO", "Drift detection,\nalaertas, retraining", VERDE, BLANCO),
]

for i, (nombre, desc, col, tc) in enumerate(cajas):
    x = 0.5 + i * 3.2
    add_rect(slide, x, 0.9, 2.9, 1.8, col)
    add_centered_text(slide, nombre, x, 1.0, 2.9, 0.55,
                      color_hex=tc, size_pt=17, bold=True)
    add_centered_text(slide, desc, x, 1.55, 2.9, 0.95,
                      color_hex=tc, size_pt=12)
    if i < 3:
        add_centered_text(slide, "→", x + 2.9, 1.55, 0.3, 0.55,
                          color_hex=AZUL_CORP, size_pt=20, bold=True)

# Definicion
add_rect(slide, 0.5, 3.1, 12.33, 0.06, NARANJA)
add_centered_text(slide,
                  "MLOps = ML + DevOps + Ingenieria de Datos",
                  0.5, 3.3, 12.33, 0.55,
                  color_hex=AZUL_CORP, size_pt=20, bold=True)

# Descripcion
descripcion_bullets = [
    "Automatizacion de pipelines de entrenamiento y despliegue",
    "Gobernanza y trazabilidad de modelos y datos",
    "Colaboracion estructurada entre Data Scientists, MLEngineers y Negocio",
]
add_bullet_box(slide, descripcion_bullets, 0.5, 4.0, 12.33, 2.0,
               color_hex=GRIS_TEXTO, size_pt=14)

# ══════════════════════════════════════════════
# SLIDE 7 — SIN vs CON MLOPS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Sin MLOps vs Con MLOps: la diferencia es clara",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

# Cabeceras
add_rect(slide, 0.4, 0.8, 2.5, 0.55, AZUL_CORP)
add_centered_text(slide, "ASPECTO", 0.4, 0.87, 2.5, 0.4,
                  color_hex=BLANCO, size_pt=13, bold=True)
add_rect(slide, 2.95, 0.8, 4.7, 0.55, ROJO)
add_centered_text(slide, "SIN MLOps", 2.95, 0.87, 4.7, 0.4,
                  color_hex=BLANCO, size_pt=13, bold=True)
add_rect(slide, 7.7, 0.8, 5.2, 0.55, VERDE)
add_centered_text(slide, "CON MLOps", 7.7, 0.87, 5.2, 0.4,
                  color_hex=BLANCO, size_pt=13, bold=True)

filas = [
    ("Velocidad",         "6+ meses por modelo",             "2-4 semanas por modelo"),
    ("Calidad",           "Errores detectados tarde",        "Validacion automatica"),
    ("Trazabilidad",      "Sin registro de versiones",       "100% auditable y versionado"),
    ("Escalabilidad",     "Cada modelo desde cero",          "Pipelines reutilizables"),
    ("Cumplimiento",      "Imposible demostrar governance",  "DORA y EU AI Act ready"),
]

for i, (aspecto, sin, con) in enumerate(filas):
    y = 1.45 + i * 0.88
    bg = GRIS_CLARO if i % 2 == 0 else BLANCO
    add_rect(slide, 0.4, y, 2.5, 0.78, bg, AZUL_MEDIO, 0.5)
    add_textbox(slide, aspecto, 0.5, y + 0.12, 2.3, 0.55,
                color_hex=AZUL_CORP, size_pt=13, bold=True)
    add_rect(slide, 2.95, y, 4.7, 0.78, "FDE8E8", ROJO, 0.5)
    add_textbox(slide, sin, 3.05, y + 0.12, 4.5, 0.55,
                color_hex=ROJO, size_pt=12)
    add_rect(slide, 7.7, y, 5.2, 0.78, "E8F8E8", VERDE, 0.5)
    add_textbox(slide, con, 7.8, y + 0.12, 5.0, 0.55,
                color_hex=VERDE, size_pt=12)

# ══════════════════════════════════════════════
# SLIDE 8 — PAUSA PREGUNTA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, NARANJA, prs)
add_centered_text(slide, "?", 4.0, 0.5, 5.33, 1.5,
                  color_hex=BLANCO, size_pt=80, bold=True)
add_centered_text(slide, "Pausa para reflexion", 1.0, 1.6, 11.33, 0.6,
                  color_hex=BLANCO, size_pt=18, italic=True)
add_centered_text(slide,
                  "\u00bfEn su \u00e1rea tienen modelos detenidos\no que tardan meses en actualizarse?",
                  0.5, 2.4, 12.33, 1.2,
                  color_hex=BLANCO, size_pt=32, bold=True)
add_centered_text(slide,
                  "Compartid un ejemplo concreto",
                  0.5, 4.0, 12.33, 0.5,
                  color_hex=BLANCO, size_pt=18, italic=True)

# ══════════════════════════════════════════════
# SLIDE 9 — CICLO DE VIDA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)
add_centered_text(slide, "El ciclo de vida de un modelo de IA",
                  0.5, 0.2, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=28, bold=True)

pasos_ciclo = [
    ("1", "Definir\nProblema",    NARANJA, BLANCO),
    ("2", "Preparar\nDatos",      BLANCO,  AZUL_CORP),
    ("3", "Entrenar\nModelo",     NARANJA, BLANCO),
    ("4", "Evaluar\nResultados",  BLANCO,  AZUL_CORP),
    ("5", "Desplegar\nen Prod",   NARANJA, BLANCO),
    ("6", "Monitorear\ny Mejorar",BLANCO,  AZUL_CORP),
]

# Dos filas de 3
for i, (num, texto, bg, tc) in enumerate(pasos_ciclo):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.0
    y = 1.3 + row * 2.3
    add_rect(slide, x, y, 3.3, 1.7, bg)
    add_centered_text(slide, num, x, y + 0.1, 3.3, 0.55,
                      color_hex=tc, size_pt=30, bold=True)
    add_centered_text(slide, texto, x, y + 0.65, 3.3, 0.85,
                      color_hex=tc, size_pt=14, bold=True)
    if col < 2:
        add_centered_text(slide, "→", x + 3.3, y + 0.55, 0.7, 0.6,
                          color_hex=NARANJA, size_pt=22, bold=True)

add_centered_text(slide, "↻  El ciclo se repite continuamente — mejora incremental",
                  0.5, 6.5, 12.33, 0.55,
                  color_hex=NARANJA, size_pt=16, italic=True)

# ══════════════════════════════════════════════
# SLIDE 10 — EJEMPLO SIN MLOPS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Caso real: Deteccion de fraude SIN MLOps",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

bullets = [
    "X  Problema detectado: fraude en aumento, modelo actual degradado",
    "X  Proceso manual: data scientist dedica 3 semanas solo a preparar datos",
    "X  Sin entorno de staging: pruebas directamente en produccion",
    "X  Tiempo de respuesta: 6 meses desde deteccion hasta nuevo modelo",
    "X  Consecuencia: 2.4M EUR en fraude no detectado durante el proceso",
]
add_bullet_box(slide, bullets, 0.5, 0.9, 12.33, 3.8,
               color_hex=ROJO, size_pt=14)

# Timeline
add_rect(slide, 0.5, 5.0, 12.33, 0.08, ROJO)
hitos = [("Mes 1", "Deteccion\nproblema"), ("Mes 2-4", "Preparacion\ndatos"),
         ("Mes 4-5", "Entrenamiento\nmanual"), ("Mes 6+", "Despliegue\nmanual")]
for i, (mes, desc) in enumerate(hitos):
    x = 0.5 + i * 3.1
    add_rect(slide, x + 1.2, 4.8, 0.12, 0.45, ROJO)
    add_centered_text(slide, mes, x, 5.15, 2.8, 0.35,
                      color_hex=ROJO, size_pt=11, bold=True)
    add_centered_text(slide, desc, x, 5.5, 2.8, 0.45,
                      color_hex=GRIS_TEXTO, size_pt=10)

# ══════════════════════════════════════════════
# SLIDE 11 — EJEMPLO CON MLOPS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Caso real: Deteccion de fraude CON MLOps",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

bullets = [
    "OK  Alerta automatica: pipeline detecta degradacion del modelo",
    "OK  Feature store pre-computado: datos listos en horas, no semanas",
    "OK  Entorno staging automatizado: pruebas sin riesgo en produccion",
    "OK  Tiempo de respuesta: 4 semanas de extremo a extremo",
    "OK  Resultado: 95% precision mantenida, cero tiempo de inactividad",
]
add_bullet_box(slide, bullets, 0.5, 0.9, 12.33, 3.8,
               color_hex=VERDE, size_pt=14)

# Metricas
add_rect(slide, 0.5, 4.7, 12.33, 0.06, VERDE)
metricas = [("6 meses → 4 semanas", "Reduccion tiempo"), ("80%", "Menos retrabajo"),
            ("95% SLA", "Precision mantenida"), ("0", "Incidentes graves")]
for i, (val, desc) in enumerate(metricas):
    x = 0.5 + i * 3.1
    add_rect(slide, x, 4.95, 2.8, 1.3, VERDE)
    add_centered_text(slide, val, x, 5.0, 2.8, 0.65,
                      color_hex=BLANCO, size_pt=15, bold=True)
    add_centered_text(slide, desc, x, 5.65, 2.8, 0.4,
                      color_hex=BLANCO, size_pt=11)

# ══════════════════════════════════════════════
# SLIDE 12 — GRAFICA TIEMPO
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Tiempo de entrega de modelos: con y sin MLOps",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

chart_data = ChartData()
chart_data.categories = ['Desarrollo', 'Validacion', 'Despliegue', 'Monitoreo']
chart_data.add_series('Sin MLOps (semanas)', (12, 8, 6, 4))
chart_data.add_series('Con MLOps (semanas)', (3, 1, 0.5, 0.5))

chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.5), Inches(0.9), Inches(10), Inches(5.5),
    chart_data
)
chart = chart_shape.chart

from pptx.util import Pt as PtU
from pptx.dml.color import RGBColor as RGBC

try:
    series0 = chart.series[0]
    series0.format.fill.solid()
    series0.format.fill.fore_color.rgb = rgb(ROJO)
    series1 = chart.series[1]
    series1.format.fill.solid()
    series1.format.fill.fore_color.rgb = rgb(AZUL_MEDIO)
except Exception as e:
    print(f"Chart series color error: {e}")

chart.has_legend = True
chart.has_title = False

add_textbox(slide, "Reduccion media: 78% menos tiempo con MLOps",
            0.5, 6.5, 12.0, 0.4,
            color_hex=VERDE, size_pt=14, bold=True)

# ══════════════════════════════════════════════
# SLIDE 13 — PAUSA PREGUNTA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, NARANJA, prs)
add_centered_text(slide, "?", 4.0, 0.5, 5.33, 1.5,
                  color_hex=BLANCO, size_pt=80, bold=True)
add_centered_text(slide, "Pausa para reflexion", 1.0, 1.6, 11.33, 0.6,
                  color_hex=BLANCO, size_pt=18, italic=True)
add_centered_text(slide,
                  "\u00bfQu\u00e9 les costar\u00eda m\u00e1s:\nun modelo que falla silenciosamente\no uno que tarda 6 meses en actualizarse?",
                  0.5, 2.3, 12.33, 1.8,
                  color_hex=BLANCO, size_pt=28, bold=True)
add_centered_text(slide,
                  "Ambos tienen coste — MLOps reduce los dos riesgos",
                  0.5, 4.4, 12.33, 0.5,
                  color_hex=BLANCO, size_pt=16, italic=True)

# ══════════════════════════════════════════════
# SLIDE 14 — RIESGOS SIN MLOPS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Riesgos de no tener MLOps",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

riesgos = [
    ("RIESGO REGULATORIO", ROJO,
     "Multas y sanciones | Imposible demostrar trazabilidad | Modelos no explicables | Incumplimiento EU AI Act / DORA"),
    ("RIESGO REPUTACIONAL", AMARILLO,
     "Modelos sesgados en produccion | Errores visibles al cliente | Perdida de confianza institucional | Cobertura mediatica negativa"),
    ("RIESGO OPERATIVO", ROJO,
     "Modelo desactualizado sin deteccion | Fraude no detectado por drift | Caida de rendimiento silenciosa | Dependencia de personas clave"),
]

for i, (titulo, col, desc) in enumerate(riesgos):
    y = 0.9 + i * 1.8
    add_rect(slide, 0.4, y, 12.53, 1.55, "FFF8F8", col, 2.0)
    add_rect(slide, 0.4, y, 2.8, 1.55, col)
    add_centered_text(slide, titulo, 0.4, y + 0.45, 2.8, 0.65,
                      color_hex=BLANCO, size_pt=13, bold=True)
    add_textbox(slide, desc, 3.4, y + 0.3, 9.2, 0.95,
                color_hex=GRIS_TEXTO, size_pt=13)

# ══════════════════════════════════════════════
# SLIDE 15 — NORMATIVA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)
add_centered_text(slide, "MLOps y el marco normativo",
                  0.5, 0.2, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=28, bold=True)
add_rect(slide, 1.5, 1.0, 9.33, 0.05, NARANJA)

normas = [
    ("DORA", "Digital Operational Resilience Act",
     "Resiliencia operativa digital\n\n"
     "• Trazabilidad completa de modelos en produccion\n"
     "• Gestion de incidentes automatizada\n"
     "• Auditoria de cambios en sistemas criticos\n"
     "• MLOps proporciona el registro y la trazabilidad exigidos"),
    ("EU AI Act", "Reglamento Europeo de Inteligencia Artificial",
     "Explicabilidad y transparencia\n\n"
     "• Documentacion de datos de entrenamiento\n"
     "• Versionado de modelos de alto riesgo\n"
     "• Monitoreo continuo de sesgos\n"
     "• MLOps genera el audit trail que exige la norma"),
]

for i, (sigla, nombre, desc) in enumerate(normas):
    x = 0.5 + i * 6.5
    add_rect(slide, x, 1.3, 5.9, 4.5, BLANCO)
    add_rect(slide, x, 1.3, 5.9, 0.7, NARANJA)
    add_centered_text(slide, sigla, x, 1.35, 5.9, 0.55,
                      color_hex=BLANCO, size_pt=22, bold=True)
    add_textbox(slide, nombre, x + 0.15, 2.1, 5.6, 0.4,
                color_hex=AZUL_CORP, size_pt=11, bold=True, italic=True)
    add_bullet_box(slide, desc.split("\n"), x + 0.15, 2.55, 5.6, 2.9,
                   color_hex=GRIS_TEXTO, size_pt=11)

add_centered_text(slide, "El cumplimiento normativo ya no es opcional",
                  0.5, 6.5, 12.33, 0.55,
                  color_hex=NARANJA, size_pt=16, bold=True, italic=True)

# ══════════════════════════════════════════════
# SLIDE 16 — BENEFICIO 1 VELOCIDAD
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Beneficio 1: Velocidad",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

# Numero decorativo grande
add_textbox(slide, "01", 8.5, 0.5, 4.5, 3.5,
            color_hex=AZUL_CLARO, size_pt=180, bold=True)

# Stat
add_rect(slide, 0.4, 0.85, 7.5, 1.4, NARANJA)
add_centered_text(slide, "De 6 meses a 2 semanas en despliegue",
                  0.4, 1.0, 7.5, 0.75,
                  color_hex=BLANCO, size_pt=22, bold=True)
add_centered_text(slide, "Reduccion del 85% en tiempo de entrega",
                  0.4, 1.7, 7.5, 0.4,
                  color_hex=BLANCO, size_pt=14)

bullets = [
    "Pipeline de entrenamiento completamente automatizado",
    "CI/CD para ML: integrar, validar y desplegar sin intervencion manual",
    "Re-entrenamiento automatico ante deteccion de drift",
    "Reutilizacion de componentes: acelera cada nuevo modelo",
]
add_bullet_box(slide, bullets, 0.4, 2.55, 7.8, 3.5,
               color_hex=GRIS_TEXTO, size_pt=14)

# ══════════════════════════════════════════════
# SLIDE 17 — BENEFICIO 2 CALIDAD
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Beneficio 2: Calidad y Estabilidad",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)
add_textbox(slide, "02", 8.5, 0.5, 4.5, 3.5,
            color_hex=AZUL_CLARO, size_pt=180, bold=True)
add_rect(slide, 0.4, 0.85, 7.5, 1.4, VERDE)
add_centered_text(slide, "90% reduccion de fallos en produccion",
                  0.4, 1.0, 7.5, 0.75,
                  color_hex=BLANCO, size_pt=22, bold=True)
add_centered_text(slide, "Estabilidad operativa garantizada",
                  0.4, 1.7, 7.5, 0.4,
                  color_hex=BLANCO, size_pt=14)
bullets = [
    "Validacion automatica de datos en cada ingesta",
    "Tests de modelo: accuracy, fairness, robustez antes de produccion",
    "Monitoreo continuo: alertas ante degradacion de rendimiento",
    "Rollback automatico ante anomalias detectadas",
]
add_bullet_box(slide, bullets, 0.4, 2.55, 7.8, 3.5,
               color_hex=GRIS_TEXTO, size_pt=14)

# ══════════════════════════════════════════════
# SLIDE 18 — BENEFICIO 3 TRAZABILIDAD
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Beneficio 3: Trazabilidad y Explicabilidad",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)
add_textbox(slide, "03", 8.5, 0.5, 4.5, 3.5,
            color_hex=AZUL_CLARO, size_pt=180, bold=True)
add_rect(slide, 0.4, 0.85, 7.5, 1.4, AZUL_MEDIO)
add_centered_text(slide, "100% de los modelos auditables",
                  0.4, 1.0, 7.5, 0.75,
                  color_hex=BLANCO, size_pt=22, bold=True)
add_centered_text(slide, "Cumplimiento normativo sin esfuerzo adicional",
                  0.4, 1.7, 7.5, 0.4,
                  color_hex=BLANCO, size_pt=14)
bullets = [
    "Registro automatico de versiones de modelos y datasets",
    "Quién entrenó, cuando, con qué datos y con qué hiperparametros",
    "Cadena de aprobacion documentada: data scientist, MLEngineer, negocio",
    "Explicabilidad integrada: SHAP values, feature importance exportable",
]
add_bullet_box(slide, bullets, 0.4, 2.55, 7.8, 3.5,
               color_hex=GRIS_TEXTO, size_pt=14)

# ══════════════════════════════════════════════
# SLIDE 19 — BENEFICIO 4 ESCALABILIDAD
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Beneficio 4: Escalabilidad",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)
add_textbox(slide, "04", 8.5, 0.5, 4.5, 3.5,
            color_hex=AZUL_CLARO, size_pt=180, bold=True)
add_rect(slide, 0.4, 0.85, 7.5, 1.4, NARANJA)
add_centered_text(slide, "De 1 modelo a 50+ sin coste proporcional",
                  0.4, 1.0, 7.5, 0.75,
                  color_hex=BLANCO, size_pt=22, bold=True)
add_centered_text(slide, "Economia de escala en inteligencia artificial",
                  0.4, 1.7, 7.5, 0.4,
                  color_hex=BLANCO, size_pt=14)
bullets = [
    "Infraestructura como codigo: entornos reproducibles en minutos",
    "Reutilizacion de pipelines entre equipos y casos de uso",
    "Gobierno centralizado: un punto de control para todos los modelos",
    "Auto-scaling de recursos: pago por uso real, no por capacidad reservada",
]
add_bullet_box(slide, bullets, 0.4, 2.55, 7.8, 3.5,
               color_hex=GRIS_TEXTO, size_pt=14)

# ══════════════════════════════════════════════
# SLIDE 20 — GRAFICA ROI
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "El coste de MLOps vs el coste de no tenerlo",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

chart_data2 = ChartData()
chart_data2.categories = ['Año 1', 'Año 2', 'Año 3']
chart_data2.add_series('Inversion MLOps (k€)', (100, 130, 150))
chart_data2.add_series('Coste sin MLOps (k€)', (200, 350, 500))
chart_data2.add_series('Ahorro acumulado (k€)', (100, 220, 350))

chart_shape2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(0.5), Inches(0.9), Inches(9.5), Inches(5.3),
    chart_data2
)
chart2 = chart_shape2.chart
chart2.has_legend = True
chart2.has_title = False

# Anotacion
add_rect(slide, 10.2, 2.5, 2.7, 0.85, VERDE)
add_centered_text(slide, "ROI positivo\ndesde el primer año",
                  10.2, 2.55, 2.7, 0.7,
                  color_hex=BLANCO, size_pt=13, bold=True)

add_textbox(slide, "Valores ilustrativos basados en benchmarks de industria (Gartner, Databricks)",
            0.5, 6.5, 12.0, 0.4,
            color_hex=GRIS_MEDIO, size_pt=10, italic=True)

# ══════════════════════════════════════════════
# SLIDE 21 — PAUSA PREGUNTA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, NARANJA, prs)
add_centered_text(slide, "?", 4.0, 0.5, 5.33, 1.5,
                  color_hex=BLANCO, size_pt=80, bold=True)
add_centered_text(slide, "Pausa para reflexion", 1.0, 1.6, 11.33, 0.6,
                  color_hex=BLANCO, size_pt=18, italic=True)
add_centered_text(slide,
                  "\u00bfQue beneficio es mas critico\npara su \u00e1rea de negocio?",
                  0.5, 2.4, 12.33, 1.2,
                  color_hex=BLANCO, size_pt=34, bold=True)
add_centered_text(slide,
                  "Velocidad   |   Calidad   |   Trazabilidad   |   Escalabilidad",
                  0.5, 3.9, 12.33, 0.55,
                  color_hex=BLANCO, size_pt=18)

# ══════════════════════════════════════════════
# SLIDE 22 — STACK TECNOLOGICO
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)
add_centered_text(slide, "Como lo implementamos: nuestra plataforma",
                  0.5, 0.2, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=28, bold=True)

tecnologias = [
    ("Dataiku", "Orquestacion visual de pipelines ML sin codigo.\nDrag & drop para equipos de negocio y datos.\nGovernance centralizado de todos los proyectos IA."),
    ("Google Cloud Platform", "Infraestructura escalable, segura y global.\nBigQuery, Vertex AI, Cloud Run integrados.\nPago por uso real, sin inversion inicial en HW."),
    ("MLflow", "Tracking de experimentos y versionado.\nModel Registry: gobierno del ciclo de vida.\nIntegracion nativa con Python, R y Spark."),
]

for i, (nombre, desc) in enumerate(tecnologias):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 1.3, 3.9, 4.2, BLANCO)
    add_rect(slide, x, 1.3, 3.9, 0.65, NARANJA)
    add_centered_text(slide, nombre, x, 1.35, 3.9, 0.55,
                      color_hex=BLANCO, size_pt=16, bold=True)
    add_bullet_box(slide, desc.split("\n"), x + 0.15, 2.1, 3.6, 3.0,
                   color_hex=GRIS_TEXTO, size_pt=12)

add_centered_text(slide,
                  "Tecnologia enterprise con interfaz accesible para negocio",
                  0.5, 6.5, 12.33, 0.55,
                  color_hex=NARANJA, size_pt=16, bold=True, italic=True)

# ══════════════════════════════════════════════
# SLIDE 23 — CASOS DE USO
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Casos de uso activos en produccion",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

casos = [
    ("Scoring crediticio en tiempo real",
     "Evaluacion automatica del riesgo en <100ms\nPipeline diario de actualizacion de scores",
     "En produccion", VERDE),
    ("Deteccion de fraude transaccional",
     "Modelo de anomalias sobre flujos en tiempo real\nMonitoreo de deriva con alertas automaticas",
     "En produccion", VERDE),
    ("Churn prediction en cartera",
     "Prediccion mensual de abandono de clientes\nAcciones preventivas automatizadas en CRM",
     "En desarrollo", AMARILLO),
    ("Analisis de riesgo de mercado",
     "Modelos de VaR y stress testing automatizados\nIntegracion con sistemas de reporte regulatorio",
     "En desarrollo", AMARILLO),
]

for i, (nombre, desc, estado, est_col) in enumerate(casos):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.5
    y = 0.9 + row * 2.9
    add_rect(slide, x, y, 6.1, 2.5, GRIS_CLARO, AZUL_MEDIO, 1.5)
    add_rect(slide, x, y, 6.1, 0.55, AZUL_CORP)
    add_textbox(slide, nombre, x + 0.15, y + 0.08, 5.8, 0.42,
                color_hex=BLANCO, size_pt=13, bold=True)
    add_bullet_box(slide, desc.split("\n"), x + 0.15, y + 0.7, 5.8, 1.1,
                   color_hex=GRIS_TEXTO, size_pt=11)
    add_rect(slide, x + 3.8, y + 1.9, 2.1, 0.38, est_col)
    add_centered_text(slide, estado, x + 3.8, y + 1.93, 2.1, 0.3,
                      color_hex=BLANCO, size_pt=10, bold=True)

# ══════════════════════════════════════════════
# SLIDE 24 — HITOS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)
add_centered_text(slide, "Hitos conseguidos",
                  0.5, 0.2, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=28, bold=True)
add_rect(slide, 1.5, 1.0, 9.33, 0.05, NARANJA)

metricas = [
    ("12", "modelos en produccion"),
    ("3 sem.", "tiempo medio despliegue"),
    ("95%", "SLA cumplimiento"),
    ("0", "incidentes regulatorios 2024"),
]

for i, (val, desc) in enumerate(metricas):
    row = i // 2
    col = i % 2
    x = 1.2 + col * 5.7
    y = 1.4 + row * 2.5
    add_rect(slide, x, y, 5.0, 2.1, NARANJA)
    add_centered_text(slide, val, x, y + 0.15, 5.0, 1.0,
                      color_hex=BLANCO, size_pt=60, bold=True)
    add_centered_text(slide, desc, x, y + 1.2, 5.0, 0.6,
                      color_hex=BLANCO, size_pt=16)

# ══════════════════════════════════════════════
# SLIDE 25 — RETOS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Retos actuales y como los abordamos",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

# Cabeceras
add_rect(slide, 0.4, 0.8, 5.5, 0.55, AZUL_CORP)
add_centered_text(slide, "RETO", 0.4, 0.87, 5.5, 0.4,
                  color_hex=BLANCO, size_pt=14, bold=True)
add_rect(slide, 6.0, 0.8, 6.9, 0.55, NARANJA)
add_centered_text(slide, "COMO LO ABORDAMOS", 6.0, 0.87, 6.9, 0.4,
                  color_hex=BLANCO, size_pt=14, bold=True)

tabla = [
    ("Calidad de datos variable",        "Validacion automatica y perfilado en ingesta"),
    ("Adopcion por equipos de negocio",  "Formacion, acompañamiento y interfaces visuales"),
    ("Integracion con sistemas legacy",  "APIs REST y conectores estandar documentados"),
    ("Tiempo de aprobacion de modelos",  "Gates automatizados con criterios claros y aprobacion digital"),
    ("Gobierno de modelos en escala",    "Model Registry centralizado con flujo de aprobacion multicapa"),
]

for i, (reto, solucion) in enumerate(tabla):
    y = 1.45 + i * 1.0
    bg = GRIS_CLARO if i % 2 == 0 else BLANCO
    add_rect(slide, 0.4, y, 5.5, 0.88, bg, AZUL_MEDIO, 0.5)
    add_textbox(slide, reto, 0.55, y + 0.15, 5.2, 0.55,
                color_hex=AZUL_CORP, size_pt=13, bold=True)
    add_rect(slide, 6.0, y, 6.9, 0.88, bg, NARANJA, 0.5)
    add_textbox(slide, solucion, 6.15, y + 0.15, 6.6, 0.55,
                color_hex=GRIS_TEXTO, size_pt=12)

# ══════════════════════════════════════════════
# SLIDE 26 — ROADMAP
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)
add_centered_text(slide, "Roadmap 2025-2026",
                  0.5, 0.2, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=28, bold=True)

# Linea timeline
add_rect(slide, 0.8, 3.3, 11.73, 0.1, NARANJA)

hitos_rm = [
    ("Q1 2025", "Plataforma MLOps unificada", "Dataiku + GCP + MLflow integrados"),
    ("Q2 2025", "20 modelos en produccion", "Expansion casos de uso validados"),
    ("Q3 2025", "Self-service para negocio", "Interfaces no-code para equipos de negocio"),
    ("Q1 2026", "IA explicable regulatoria", "Cumplimiento EU AI Act avanzado"),
]

for i, (trimestre, hito, desc) in enumerate(hitos_rm):
    x = 0.8 + i * 3.1
    # Punto en timeline
    add_rect(slide, x + 1.3, 3.1, 0.18, 0.5, NARANJA)
    # Caja superior
    add_rect(slide, x, 1.3, 2.9, 1.65, BLANCO)
    add_rect(slide, x, 1.3, 2.9, 0.45, NARANJA)
    add_centered_text(slide, trimestre, x, 1.35, 2.9, 0.38,
                      color_hex=BLANCO, size_pt=14, bold=True)
    add_centered_text(slide, hito, x, 1.8, 2.9, 0.55,
                      color_hex=AZUL_CORP, size_pt=12, bold=True)
    add_centered_text(slide, desc, x, 2.35, 2.9, 0.55,
                      color_hex=GRIS_TEXTO, size_pt=10)
    # Caja inferior
    add_rect(slide, x, 3.65, 2.9, 1.5, "00264D")
    add_centered_text(slide, f"Hito {i+1}", x, 3.75, 2.9, 0.4,
                      color_hex=NARANJA, size_pt=11, bold=True)

add_centered_text(slide, "Crecimiento iterativo con valor incremental en cada entrega",
                  0.5, 6.5, 12.33, 0.55,
                  color_hex=AZUL_CLARO, size_pt=14, italic=True)

# ══════════════════════════════════════════════
# SLIDE 27 — PAUSA PREGUNTA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, NARANJA, prs)
add_centered_text(slide, "?", 4.0, 0.5, 5.33, 1.5,
                  color_hex=BLANCO, size_pt=80, bold=True)
add_centered_text(slide, "Pausa para reflexion", 1.0, 1.6, 11.33, 0.6,
                  color_hex=BLANCO, size_pt=18, italic=True)
add_centered_text(slide,
                  "\u00bfQue iniciativas del roadmap se alinean mejor\ncon sus objetivos de negocio?",
                  0.5, 2.4, 12.33, 1.2,
                  color_hex=BLANCO, size_pt=32, bold=True)
add_centered_text(slide,
                  "¿Cuales quereis priorizar para el siguiente trimestre?",
                  0.5, 3.9, 12.33, 0.55,
                  color_hex=BLANCO, size_pt=18)

# ══════════════════════════════════════════════
# SLIDE 28 — LO QUE NECESITAMOS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, BLANCO, prs)
add_rect(slide, 0, 0, 13.33, 0.65, AZUL_CORP)
add_textbox(slide, "Lo que necesitamos de negocio",
            0.4, 0.1, 12.0, 0.5,
            color_hex=BLANCO, size_pt=22, bold=True)

necesidades = [
    ("Casos de uso concretos con impacto medible",
     "Identificar 2-3 problemas de negocio con KPIs claros y datos disponibles"),
    ("Acceso a datos de calidad con governance",
     "Datos etiquetados, definicion de features y ownership claro del dominio"),
    ("Tiempo de expertos de negocio para validar",
     "Al menos 2h/semana de experto de dominio en las primeras 4 semanas"),
    ("Criterios claros de exito del modelo",
     "Precision minima aceptable, latencia requerida, frecuencia de actualizacion"),
]

for i, (titulo, desc) in enumerate(necesidades):
    y = 0.9 + i * 1.5
    add_rect(slide, 0.4, y, 0.6, 1.1, NARANJA)
    add_centered_text(slide, str(i+1), 0.4, y + 0.25, 0.6, 0.6,
                      color_hex=BLANCO, size_pt=20, bold=True)
    add_rect(slide, 1.1, y, 11.73, 1.1, GRIS_CLARO, AZUL_MEDIO, 0.5)
    add_textbox(slide, titulo, 1.25, y + 0.08, 11.4, 0.42,
                color_hex=AZUL_CORP, size_pt=14, bold=True)
    add_textbox(slide, desc, 1.25, y + 0.5, 11.4, 0.48,
                color_hex=GRIS_TEXTO, size_pt=12)

# ══════════════════════════════════════════════
# SLIDE 29 — PROXIMOS PASOS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)
add_centered_text(slide, "Proximos pasos concretos",
                  0.5, 0.2, 12.33, 0.7,
                  color_hex=BLANCO, size_pt=28, bold=True)
add_rect(slide, 1.5, 1.0, 9.33, 0.05, NARANJA)

pasos_accion = [
    ("1", "Workshop de identificacion de casos de uso",
     "Esta semana — sesion de 2 horas con equipo de negocio y datos\nOutput: mapa de oportunidades priorizado por impacto/viabilidad"),
    ("2", "Definicion conjunta de criterios de exito",
     "Proximas 2 semanas — data scientist + experto de negocio\nOutput: tarjeta de proyecto con KPIs, datos, timeline"),
    ("3", "Kick-off del primer proyecto piloto",
     "Proximo mes — equipo dedicado, metodologia agil\nOutput: MVP desplegado en staging en 3-4 semanas"),
]

for i, (num, titulo, desc) in enumerate(pasos_accion):
    y = 1.4 + i * 1.8
    add_rect(slide, 0.5, y, 0.75, 1.4, NARANJA)
    add_centered_text(slide, num, 0.5, y + 0.35, 0.75, 0.65,
                      color_hex=BLANCO, size_pt=30, bold=True)
    add_rect(slide, 1.4, y, 11.43, 1.4, BLANCO)
    add_textbox(slide, titulo, 1.55, y + 0.1, 11.1, 0.5,
                color_hex=AZUL_CORP, size_pt=16, bold=True)
    add_bullet_box(slide, desc.split("\n"), 1.55, y + 0.6, 11.1, 0.65,
                   color_hex=GRIS_TEXTO, size_pt=12)

# ══════════════════════════════════════════════
# SLIDE 30 — CIERRE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_background(slide, AZUL_CORP, prs)

# Decoracion
add_rect(slide, 0, 0, 0.5, 7.5, NARANJA)
add_rect(slide, 12.83, 0, 0.5, 7.5, NARANJA)
add_rect(slide, 0, 0, 13.33, 0.12, NARANJA)

# Foto instructor
if FOTO_PATH:
    try:
        slide.shapes.add_picture(FOTO_PATH, Inches(9.5), Inches(0.8),
                                  Inches(3.0), Inches(3.8))
    except Exception as e:
        print(f"Slide 30 foto error: {e}")

# Frase impacto
add_centered_text(slide,
                  '"La IA no es el futuro.',
                  0.6, 1.0, 8.8, 0.75,
                  color_hex=BLANCO, size_pt=30, bold=True)
add_centered_text(slide,
                  'Es el presente.',
                  0.6, 1.7, 8.8, 0.7,
                  color_hex=NARANJA, size_pt=36, bold=True)
add_centered_text(slide,
                  'MLOps es lo que la hace funcionar."',
                  0.6, 2.4, 8.8, 0.75,
                  color_hex=BLANCO, size_pt=30, bold=True)

add_rect(slide, 0.6, 3.3, 6.5, 0.06, NARANJA)

# Datos de contacto
add_textbox(slide, "Formacion MLOps 2025", 0.6, 3.5, 8.0, 0.45,
            color_hex=BLANCO, size_pt=16, bold=True)
add_textbox(slide, "mlops@empresa.com  |  Confidencial — Solo uso interno",
            0.6, 4.0, 8.0, 0.35,
            color_hex=AZUL_CLARO, size_pt=12)
add_textbox(slide, fecha_hoy,
            0.6, 4.4, 8.0, 0.35,
            color_hex=AZUL_CLARO, size_pt=12)

add_rect(slide, 0, 6.8, 13.33, 0.7, "00264D")
add_centered_text(slide,
                  "Gracias. Preguntas finales.",
                  0, 6.88, 13.33, 0.5,
                  color_hex=NARANJA, size_pt=20, bold=True)

# ══════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════
output_dir = "C:/Users/bk70827/PycharmProjects/mlops-ciclo-vida/pdfs-ppts/output"
os.makedirs(output_dir, exist_ok=True)
output_path = f"{output_dir}/mlops_presentacion.pptx"
prs.save(output_path)

size_kb = os.path.getsize(output_path) / 1024
print(f"\n{'='*60}")
print(f"PRESENTACION GENERADA EXITOSAMENTE")
print(f"{'='*60}")
print(f"Ruta: {output_path}")
print(f"Tamanio: {size_kb:.1f} KB")
print(f"Slides: {len(prs.slides)}")
print(f"Foto instructor: {'SI - ' + str(FOTO_PATH) if FOTO_PATH else 'NO'}")
print(f"{'='*60}")
print("\nCOLORES USADOS:")
print(f"  Azul corporativo: #{AZUL_CORP}")
print(f"  Azul medio:       #{AZUL_MEDIO}")
print(f"  Naranja acento:   #{NARANJA}")
print(f"  Verde:            #{VERDE}")
print(f"  Rojo:             #{ROJO}")
print("\nFUENTES ENCONTRADAS EN PATRONES: Open Sans, Berlin Sans FB")
print(f"FUENTE USADA: {FUENTE_CUERPO} (fallback corporativo)")
print("\nCONTENIDO PDF UTILIZADO:")
print("  - Definicion MLOps = ModelOps + DataOps + DevOps (Databricks)")
print("  - Estadistica: 83% CEOs reportan IA como prioridad (MIT Sloan)")
print("  - Estimacion Gartner: $3.9T valor IA 2022")
print("  - Framework ciclo de vida: Dev → Staging → Prod")

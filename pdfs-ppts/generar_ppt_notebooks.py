"""
Generador de presentación MLOps didáctica con imágenes reales de los notebooks.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── COLORES CORPORATIVOS ───────────────────────────────────────────────────
AZUL       = RGBColor(0x00, 0x3A, 0x70)
NARANJA    = RGBColor(0xFF, 0x6B, 0x00)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
GRIS       = RGBColor(0xF5, 0xF5, 0xF5)
GRIS_TEXTO = RGBColor(0x33, 0x33, 0x33)
VERDE      = RGBColor(0x27, 0xAE, 0x60)
AZUL_CLARO = RGBColor(0x00, 0x6E, 0xC7)
AMARILLO   = RGBColor(0xFF, 0xF0, 0xCC)
NARANJA_CLARO = RGBColor(0xFF, 0xE5, 0xCC)

# ─── RUTAS ─────────────────────────────────────────────────────────────────
BASE       = Path("C:/Users/bk70827/PycharmProjects/mlops-ciclo-vida/pdfs-ppts")
ASSETS     = BASE / "assets"
NB_IMAGES  = ASSETS / "nb_images"
OUTPUT     = BASE / "output" / "mlops_ejemplo_notebooks.pptx"
FOTO_INS   = ASSETS / "foto_instructor.png"

def img(nombre):
    return str(NB_IMAGES / nombre)

# ─── HELPERS ────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, border_color=None, alpha=None):
    shape = slide.shapes.add_shape(
        1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, texto, x, y, w, h, color, size_pt,
             bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.color.rgb = color
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox


def add_text_multiline(slide, lineas, x, y, w, h, color, size_pt,
                       bold=False, align=PP_ALIGN.LEFT, line_spacing=1.15):
    """Añade texto con múltiples párrafos."""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    first = True
    for linea in lineas:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = linea
        run.font.color.rgb = color
        run.font.size = _Pt(size_pt)
        run.font.bold = bold
        run.font.name = "Calibri"
    return txBox


def add_image_safe(slide, img_path, x, y, w, h):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
        return True
    else:
        shape = add_rect(slide, x, y, w, h, RGBColor(0xCC, 0xCC, 0xCC))
        add_text(slide, f"[Imagen: {os.path.basename(img_path)}]",
                 x + 0.1, y + h / 2 - 0.2, w - 0.2, 0.4,
                 GRIS_TEXTO, 10, align=PP_ALIGN.CENTER)
        return False


def slide_fondo(slide, color):
    add_rect(slide, 0, 0, 13.33, 7.5, color)


def slide_barra_top(slide, barra_color=AZUL, linea_color=NARANJA):
    add_rect(slide, 0, 0, 13.33, 0.55, barra_color)
    add_rect(slide, 0, 0.55, 13.33, 0.06, linea_color)


def slide_barra_lateral(slide):
    add_rect(slide, 0, 0, 0.35, 7.5, AZUL)
    add_rect(slide, 0.35, 0, 0.06, 7.5, NARANJA)


def add_numero_slide(slide, numero, total=25):
    add_text(slide, f"{numero} / {total}",
             12.3, 7.1, 1.0, 0.35,
             RGBColor(0xAA, 0xAA, 0xAA), 9, align=PP_ALIGN.RIGHT)


def add_bullet_box(slide, bullets, x, y, w, h,
                   bg_color=GRIS, text_color=GRIS_TEXTO,
                   bullet_char="▶", size_pt=13):
    add_rect(slide, x, y, w, h, bg_color, border_color=RGBColor(0xDD, 0xDD, 0xDD))
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.12),
        Inches(w - 0.3), Inches(h - 0.24)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{bullet_char}  {b}"
        run.font.color.rgb = text_color
        run.font.size = Pt(size_pt)
        run.font.name = "Calibri"


def add_caja_kpi(slide, titulo, valor, x, y, w, h,
                 bg=AZUL, text_color=BLANCO, val_color=NARANJA):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, titulo, x + 0.1, y + 0.08, w - 0.2, 0.35,
             text_color, 11, align=PP_ALIGN.CENTER)
    add_text(slide, valor, x + 0.1, y + 0.42, w - 0.2, 0.5,
             val_color, 22, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  CONSTRUCCIÓN DE LA PRESENTACIÓN
# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completamente en blanco

slide_titles = []
img_inserted = 0
img_missing  = 0

def new_slide(title_log=""):
    slide = prs.slides.add_slide(blank_layout)
    slide_titles.append(title_log)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("PORTADA — MLOps en Acción")
slide_fondo(s, AZUL)

# Franja naranja central decorativa
add_rect(s, 0, 3.1, 13.33, 0.07, NARANJA)

# Línea naranja superior
add_rect(s, 0.8, 1.55, 11.73, 0.06, NARANJA)

# Título grande
add_text(s, "MLOps en Acción",
         1.0, 1.7, 10.5, 1.4,
         BLANCO, 54, bold=True, align=PP_ALIGN.CENTER)

# Subtítulo
add_text(s, "Predicción de Precio de Viviendas",
         1.5, 3.25, 10.3, 0.65,
         NARANJA, 26, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Ejemplo paso a paso — Dataset California Housing",
         1.5, 3.9, 10.3, 0.5,
         BLANCO, 18, align=PP_ALIGN.CENTER)

# Foto instructor (lado derecho)
foto_path = str(FOTO_INS)
if os.path.exists(foto_path):
    s.shapes.add_picture(foto_path, Inches(10.5), Inches(4.6), Inches(2.0), Inches(2.0))
    add_text(s, "Instructor", 10.5, 6.65, 2.0, 0.35, BLANCO, 11, align=PP_ALIGN.CENTER)
    img_inserted += 1

# Tags
add_text(s, "🎓  Formación MLOps  ·  6 Notebooks  ·  Código reproducible",
         1.5, 6.7, 10.3, 0.45,
         RGBColor(0xAA, 0xCC, 0xFF), 13, align=PP_ALIGN.CENTER)

add_numero_slide(s, 1)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — ¿Con qué datos trabajamos?
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Dataset: California Housing")
slide_fondo(s, BLANCO)
slide_barra_lateral(s)

# Título
add_text(s, "Dataset: California Housing",
         0.6, 0.18, 11.0, 0.65,
         AZUL, 30, bold=True)

add_rect(s, 0.6, 0.85, 5.5, 0.04, NARANJA)

# Texto descriptivo izquierda
add_text_multiline(s,
    ["20.640 viviendas de California  (1990 Census)",
     "",
     "Objetivo: predecir el precio mediano de una vivienda.",
     "",
     "8 variables de entrada:",
     "  • MedInc — ingreso mediano del vecindario",
     "  • HouseAge — antigüedad de las viviendas",
     "  • AveRooms — habitaciones por hogar",
     "  • AveBedrms — dormitorios por hogar",
     "  • Population — población del bloque",
     "  • AveOccup — personas por hogar",
     "  • Latitude / Longitude — ubicación geográfica"],
    0.6, 1.0, 5.6, 5.2,
    GRIS_TEXTO, 13)

# Imagen EDA img00
r = add_image_safe(s, img("01_eda_exploratorio_img00.png"), 6.4, 0.85, 6.7, 5.6)
if r: img_inserted += 1
else: img_missing += 1

# Pie de página
add_rect(s, 0.35, 7.0, 12.98, 0.4, GRIS)
add_text(s, "Fuente: sklearn.datasets — datos reales del censo de California 1990",
         0.6, 7.05, 10.0, 0.35,
         RGBColor(0x88, 0x88, 0x88), 10, italic=True)
add_numero_slide(s, 2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Análisis Exploratorio: distribución y heatmap
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("¿Qué nos dicen los datos? Análisis Exploratorio")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "¿Qué nos dicen los datos?  Análisis Exploratorio",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

# Imagen izquierda
r = add_image_safe(s, img("01_eda_exploratorio_img01.png"), 0.4, 1.45, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

# Texto derecha
add_text(s, "¿Qué observamos?", 8.5, 1.45, 4.5, 0.45, AZUL, 16, bold=True)
add_bullet_box(s,
    ["La variable objetivo (precio) tiene distribución sesgada a la derecha — hay muchas viviendas baratas y pocas muy caras.",
     "Las variables de ingresos y precio están correlacionadas positivamente.",
     "Existen valores techo en 5.0 ($500k) que indican truncamiento artificial en los datos.",
     "Los datos geográficos permiten segmentar por zonas costeras vs interiores."],
    8.5, 1.95, 4.6, 4.5,
    bg_color=RGBColor(0xE8, 0xF0, 0xFE), text_color=GRIS_TEXTO,
    bullet_char="◉", size_pt=12)

add_numero_slide(s, 3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — EDA img02: Correlación entre variables
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Correlación entre variables — ¿qué predice el precio?")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Correlación entre variables — ¿qué predice el precio?",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

# Imagen centrada + texto abajo
r = add_image_safe(s, img("01_eda_exploratorio_img02.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "Decisiones clave:", 8.5, 1.45, 4.5, 0.4, AZUL, 16, bold=True)
add_bullet_box(s,
    ["MedInc (ingresos) es la variable con mayor correlación con el precio (r ≈ 0.69). Es el predictor más potente.",
     "Latitude correlaciona negativamente: al norte los precios bajan.",
     "AveRooms tiene correlación moderada — más habitaciones = precio mayor.",
     "Pocas correlaciones entre predictores → no hay multicolinealidad severa."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=NARANJA_CLARO, text_color=GRIS_TEXTO,
    bullet_char="▶", size_pt=12)

add_numero_slide(s, 4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — EDA img03: Distribución geográfica
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Distribución geográfica del precio en California")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Distribución geográfica del precio en California",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

r = add_image_safe(s, img("01_eda_exploratorio_img03.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "Insights geográficos:", 8.5, 1.45, 4.5, 0.4, AZUL, 16, bold=True)
add_bullet_box(s,
    ["La costa de California concentra los precios más altos (San Francisco, Los Ángeles).",
     "Las zonas del interior (Valle Central) tienen precios significativamente menores.",
     "La latitud y longitud son features relevantes — el modelo puede aprender patrones geográficos.",
     "Clustering visual: zonas metropolitanas son outliers de precio hacia arriba."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=RGBColor(0xE8, 0xF8, 0xEE), text_color=GRIS_TEXTO,
    bullet_char="▶", size_pt=12)

add_numero_slide(s, 5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — EDA img04: Variables más influyentes
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Variables más influyentes en el precio")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Variables más influyentes en el precio",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

r = add_image_safe(s, img("01_eda_exploratorio_img04.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "¿Qué variables usar?", 8.5, 1.45, 4.5, 0.4, AZUL, 16, bold=True)
add_bullet_box(s,
    ["Las visualizaciones de distribución revelan la forma real de cada variable.",
     "Variables con alta asimetría requieren transformación (log) antes del modelado.",
     "Population y AveOccup muestran outliers extremos que deben tratarse.",
     "Este análisis guía directamente la etapa de Ingeniería de Features."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=RGBColor(0xE8, 0xF0, 0xFE), text_color=GRIS_TEXTO,
    bullet_char="▶", size_pt=12)

add_numero_slide(s, 6)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — EDA img05: Mapa de calor / distribuciones avanzadas
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Outliers y patrones identificados en el EDA")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Outliers y patrones identificados en el EDA",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

r = add_image_safe(s, img("01_eda_exploratorio_img05.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "Hallazgos:", 8.5, 1.45, 4.5, 0.4, AZUL, 16, bold=True)
add_bullet_box(s,
    ["Los scatter plots revelan relaciones no lineales que los modelos lineales simples no pueden capturar.",
     "Existe heteroscedasticidad: la varianza del precio aumenta con los ingresos.",
     "Identificamos grupos de puntos con patrones distintos → posibles segmentos de mercado.",
     "Conclusión: necesitamos modelos no lineales (Random Forest, Gradient Boosting)."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=NARANJA_CLARO, text_color=GRIS_TEXTO,
    bullet_char="▶", size_pt=12)

add_numero_slide(s, 7)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — EDA img06: Relación ingresos-precio
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Relación entre ingresos y precio de vivienda")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Relación entre ingresos y precio de vivienda",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

r = add_image_safe(s, img("01_eda_exploratorio_img06.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "La variable clave:", 8.5, 1.45, 4.5, 0.4, AZUL, 16, bold=True)
add_bullet_box(s,
    ["MedInc (ingreso mediano del vecindario) explica casi el 50% de la varianza del precio.",
     "La relación es positiva pero no perfectamente lineal — a mayores ingresos, la dispersión crece.",
     "El techo artificial en $500k se ve claramente: muchos puntos en esa línea horizontal.",
     "Para el negocio: predecir precios en vecindarios ricos es más incierto que en zonas de ingreso medio."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=RGBColor(0xE8, 0xF8, 0xEE), text_color=GRIS_TEXTO,
    bullet_char="▶", size_pt=12)

add_numero_slide(s, 8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Ingeniería de Features (fondo azul)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Etapa 3: Ingeniería de Features")
slide_fondo(s, AZUL)

# Línea naranja decorativa
add_rect(s, 0.8, 0.6, 11.73, 0.06, NARANJA)

add_text(s, "Etapa 3: Ingeniería de Features",
         0.8, 0.7, 11.0, 0.75,
         BLANCO, 32, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "De datos crudos a variables inteligentes",
         0.8, 1.45, 11.0, 0.5,
         NARANJA, 20, align=PP_ALIGN.CENTER)

# 5 features derivadas en cajas
features = [
    ("rooms_per_person",    "AveRooms / AveOccup",         "Mide hacinamiento real por persona"),
    ("bedrooms_ratio",      "AveBedrms / AveRooms",        "Proporción dormitorios vs habitaciones"),
    ("income_per_room",     "MedInc / AveRooms",           "Poder adquisitivo relativo al tamaño"),
    ("population_density",  "Population / AveOccup",       "Densidad del bloque censal"),
    ("location_cluster",    "Lat × Lon (binning)",         "Zona geográfica categorizada"),
]

fx, fy, fw, fh = 0.5, 2.0, 3.8, 0.85
for i, (feat, formula, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    cx = 0.5 + col * 4.27
    cy = 2.05 + row * 1.1
    add_rect(s, cx, cy, 4.1, 0.92, BLANCO)
    add_text(s, feat, cx + 0.15, cy + 0.05, 3.8, 0.35,
             AZUL, 13, bold=True)
    add_text(s, formula, cx + 0.15, cy + 0.38, 3.8, 0.27,
             NARANJA, 11)
    add_text(s, desc, cx + 0.15, cy + 0.62, 3.8, 0.27,
             GRIS_TEXTO, 10)

# Imagen features
r = add_image_safe(s, img("02_ingenieria_features_img00.png"), 9.0, 1.9, 4.0, 3.5)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "8 variables originales  →  13 variables enriquecidas",
         0.8, 5.5, 11.0, 0.5,
         NARANJA, 16, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Más features relevantes = modelos más precisos",
         0.8, 5.95, 11.0, 0.4,
         BLANCO, 13, align=PP_ALIGN.CENTER)

add_numero_slide(s, 9)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Data Leakage y Escalado
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("¿Por qué es crítico el escalado? Evitando el Data Leakage")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "¿Por qué es crítico el escalado? Evitando el Data Leakage",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 22, bold=True)

# Caja concepto clave naranja
add_rect(s, 0.4, 1.42, 5.8, 1.0, NARANJA)
add_text(s, "⚠  DATA LEAKAGE",
         0.55, 1.5, 5.5, 0.4,
         BLANCO, 18, bold=True)
add_text(s, "Contaminar el futuro con información del presente",
         0.55, 1.85, 5.5, 0.5,
         BLANCO, 13)

# Regla MLOps
add_rect(s, 0.4, 2.55, 5.8, 1.35, RGBColor(0x00, 0x5A, 0xA0))
add_text(s, "Regla MLOps — Escalado correcto:",
         0.55, 2.62, 5.5, 0.4,
         NARANJA, 13, bold=True)
add_text_multiline(s,
    ["✅  scaler.fit_transform( X_TRAIN )   ← correcto",
     "✅  scaler.transform( X_TEST )        ← correcto",
     "❌  scaler.fit_transform( X_ALL )     ← DATA LEAKAGE"],
    0.55, 3.0, 5.5, 0.85,
    BLANCO, 11)

# Imagen
r = add_image_safe(s, img("02_ingenieria_features_img01.png"), 6.5, 1.42, 6.6, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "El scaler se ajusta SOLO con datos de entrenamiento. Nunca con datos de test.",
         0.4, 3.95, 5.8, 0.7,
         GRIS_TEXTO, 12, italic=True)

add_text(s, "Si ajustamos con todo el dataset, el modelo 've' el futuro durante el entrenamiento → métricas infladas falsamente.",
         0.4, 4.7, 5.8, 0.9,
         GRIS_TEXTO, 12)

add_numero_slide(s, 10)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Entrenamiento + MLflow (fondo azul)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Etapa 4: Entrenamiento + Tracking con MLflow")
slide_fondo(s, AZUL)
add_rect(s, 0.8, 0.6, 11.73, 0.06, NARANJA)

add_text(s, "Entrenamos 3 algoritmos y registramos todo",
         0.8, 0.7, 11.0, 0.75,
         BLANCO, 30, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Etapa 4: Entrenamiento + Tracking con MLflow",
         0.8, 1.45, 11.0, 0.45,
         NARANJA, 17, align=PP_ALIGN.CENTER)

# 3 bloques algoritmos
algos = [
    ("Regresión Lineal", "Simple e interpretable.\nBase de comparación.\nRápida (< 1 seg).", "Baseline"),
    ("Random Forest", "Potente, maneja no-linealidades.\nRobusta ante outliers.\n45 seg entrenamiento.", "Competidor"),
    ("Gradient Boosting", "El más preciso.\nOptimiza iterativamente.\n38 seg entrenamiento.", "★  Ganador"),
]

for i, (nombre, desc, tag) in enumerate(algos):
    cx = 0.6 + i * 4.2
    add_rect(s, cx, 2.1, 3.9, 3.2, BLANCO)
    # tag
    tag_color = NARANJA if "Ganador" in tag else AZUL_CLARO
    add_rect(s, cx, 2.1, 3.9, 0.45, tag_color)
    add_text(s, tag, cx + 0.1, 2.13, 3.7, 0.38,
             BLANCO, 13, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, nombre, cx + 0.1, 2.6, 3.7, 0.5,
             AZUL, 16, bold=True, align=PP_ALIGN.CENTER)
    add_text_multiline(s, desc.split("\n"),
                       cx + 0.15, 3.15, 3.6, 1.8,
                       GRIS_TEXTO, 12)

# Imagen MLflow
r = add_image_safe(s, img("03_entrenamiento_mlflow_img00.png"), 0.5, 5.35, 8.5, 1.85)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "MLflow registra: parámetros · métricas · artefactos · versión del modelo",
         9.2, 5.7, 3.8, 0.8,
         NARANJA, 13, bold=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 11)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Comparativa de experimentos (tabla)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Comparativa de experimentos en MLflow")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Comparativa de experimentos en MLflow",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

# Tabla manual
headers = ["Algoritmo", "RMSE ↓", "R² ↑", "Tiempo", "Resultado"]
rows = [
    ("Regresión Lineal",  "0.73", "0.62", "1 seg",  "Descartado"),
    ("Random Forest",     "0.52", "0.79", "45 seg", "Alternativa"),
    ("Gradient Boosting", "0.43", "0.84", "38 seg", "✅  ELEGIDO"),
]
col_widths = [3.2, 1.8, 1.5, 1.8, 2.4]
col_starts = [0.5]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

# Header
ty = 1.5
for j, (h, x, w) in enumerate(zip(headers, col_starts, col_widths)):
    add_rect(s, x, ty, w - 0.05, 0.55, AZUL)
    add_text(s, h, x + 0.1, ty + 0.08, w - 0.2, 0.4,
             BLANCO, 14, bold=True, align=PP_ALIGN.CENTER)

# Filas
for i, row in enumerate(rows):
    ty = 2.1 + i * 0.75
    is_winner = i == 2
    bg = NARANJA if is_winner else (GRIS if i % 2 == 0 else BLANCO)
    tc = BLANCO if is_winner else GRIS_TEXTO
    for j, (cell, x, w) in enumerate(zip(row, col_starts, col_widths)):
        add_rect(s, x, ty, w - 0.05, 0.65, bg,
                 border_color=RGBColor(0xDD, 0xDD, 0xDD))
        add_text(s, cell, x + 0.1, ty + 0.15, w - 0.2, 0.4,
                 tc, 13, bold=is_winner, align=PP_ALIGN.CENTER)

# Explicación
add_rect(s, 0.5, 4.45, 10.7, 1.35, RGBColor(0xE8, 0xF0, 0xFE),
         border_color=AZUL_CLARO)
add_text(s, "¿Por qué Gradient Boosting?",
         0.7, 4.55, 5.0, 0.4, AZUL, 14, bold=True)
add_text_multiline(s,
    ["• Mejor RMSE (0.43 vs 0.52): las predicciones se desvían en promedio $43k — dentro del umbral aceptable.",
     "• Mejor R² (0.84): explica el 84% de la varianza del precio.",
     "• El tiempo de entrenamiento es similar a Random Forest pero con mejor resultado."],
    0.7, 4.95, 10.3, 0.75,
    GRIS_TEXTO, 12)

add_numero_slide(s, 12)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Validación cruzada
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Validación cruzada — resultados robustos")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Etapa 4: Selección rigurosa del modelo",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

add_text(s, "Validación cruzada — resultados robustos",
         0.4, 1.25, 8.0, 0.45,
         NARANJA, 18, bold=True)

# Imagen
r = add_image_safe(s, img("04_seleccion_modelo_img00.png"), 0.4, 1.75, 8.0, 4.9)
if r: img_inserted += 1
else: img_missing += 1

# Panel explicativo derecha
add_text(s, "¿Qué es K-Fold?", 8.7, 1.45, 4.3, 0.4, AZUL, 15, bold=True)
add_bullet_box(s,
    ["Dividimos los datos en 5 subconjuntos (folds).",
     "Entrenamos 5 veces, cada vez con un fold distinto como test.",
     "El resultado final es el promedio de las 5 evaluaciones.",
     "Garantiza que el resultado no es casualidad de una sola partición.",
     "Si los 5 resultados son similares → modelo estable y generalizable."],
    8.7, 1.9, 4.3, 3.5,
    bg_color=RGBColor(0xE8, 0xF8, 0xEE), text_color=GRIS_TEXTO,
    bullet_char="◉", size_pt=12)

add_rect(s, 8.7, 5.5, 4.3, 0.95, AZUL)
add_text(s, "\"5-fold CV garantiza que el resultado no es suerte\"",
         8.85, 5.6, 4.0, 0.75,
         NARANJA, 13, italic=True)

add_numero_slide(s, 13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Curvas de aprendizaje
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Curvas de aprendizaje — ¿overfitting o underfitting?")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Curvas de aprendizaje — ¿overfitting o underfitting?",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

# Imagen
r = add_image_safe(s, img("04_seleccion_modelo_img01.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

# Panel derecha — bias-variance
add_text(s, "Bias-Variance Tradeoff:", 8.5, 1.45, 4.5, 0.4, AZUL, 15, bold=True)

cuadros = [
    (NARANJA,   "Sobreajuste (Overfitting)",
     "El modelo memoriza el dataset de entrenamiento pero falla en datos nuevos. Train error ↓↓  /  Val error ↑↑"),
    (RGBColor(0xAA, 0xAA, 0xAA), "Subajuste (Underfitting)",
     "El modelo es demasiado simple para capturar los patrones. Ambos errores son altos."),
    (VERDE,     "Nuestro modelo: equilibrio correcto ✅",
     "Las curvas de entrenamiento y validación convergen. El error es bajo y consistente en ambos conjuntos."),
]

cy = 1.95
for bg, titulo, desc in cuadros:
    add_rect(s, 8.5, cy, 4.6, 1.25, bg)
    add_text(s, titulo, 8.65, cy + 0.07, 4.3, 0.38,
             BLANCO, 13, bold=True)
    add_text(s, desc, 8.65, cy + 0.45, 4.3, 0.7,
             BLANCO, 11)
    cy += 1.4

add_numero_slide(s, 14)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Selección img02 adicional
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Análisis de residuos — diagnóstico del modelo")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Análisis de residuos — diagnóstico del modelo",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 24, bold=True)

r = add_image_safe(s, img("04_seleccion_modelo_img02.png"), 1.5, 1.42, 10.0, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "Un buen modelo tiene residuos distribuidos aleatoriamente alrededor de cero, sin patrones sistemáticos.",
         1.5, 6.82, 10.0, 0.55,
         RGBColor(0x88, 0x88, 0x88), 11, italic=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 15)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Evaluación Final (fondo azul)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Etapa 5: Evaluación Final — El modelo se evalúa con datos NUNCA vistos")
slide_fondo(s, AZUL)
add_rect(s, 0.8, 0.6, 11.73, 0.06, NARANJA)

add_text(s, "Etapa 5: Evaluación Final",
         0.8, 0.7, 11.0, 0.75,
         BLANCO, 32, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "El modelo se evalúa con datos NUNCA vistos durante el entrenamiento",
         0.8, 1.45, 11.0, 0.5,
         NARANJA, 17, align=PP_ALIGN.CENTER)

# Concepto clave
add_rect(s, 0.7, 2.1, 11.93, 0.85, RGBColor(0x00, 0x2A, 0x55))
add_text(s, "El 20% de los datos se reserva desde el inicio. El modelo nunca los ve durante el entrenamiento.",
         0.9, 2.2, 11.5, 0.6,
         NARANJA, 15, align=PP_ALIGN.CENTER)

# KPIs
kpis = [
    ("RMSE",  "0.43",  "Umbral: < 0.50  ✅", VERDE),
    ("R²",    "0.84",  "Umbral: > 0.80  ✅", VERDE),
    ("MAE",   "0.31",  "Error medio abs  ✅", VERDE),
    ("MAPE",  "16.8%", "Error porcentual  ✅", VERDE),
]
for i, (label, val, umbral, col) in enumerate(kpis):
    cx = 1.0 + i * 2.85
    add_rect(s, cx, 3.15, 2.6, 1.55, BLANCO)
    add_text(s, label,  cx + 0.1, 3.22, 2.4, 0.4,
             AZUL, 16, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, val,    cx + 0.1, 3.6,  2.4, 0.55,
             col, 28, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, umbral, cx + 0.1, 4.12, 2.4, 0.45,
             GRIS_TEXTO, 11, align=PP_ALIGN.CENTER)

# Veredicto
add_rect(s, 2.5, 5.0, 8.33, 1.1, VERDE)
add_text(s, "✅  MODELO APROBADO — LISTO PARA PRODUCCIÓN",
         2.6, 5.2, 8.1, 0.7,
         BLANCO, 24, bold=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 16)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — Predicciones vs Realidad
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Predicciones vs Realidad")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Predicciones vs Realidad",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

r = add_image_safe(s, img("05_evaluacion_final_img00.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "Interpretación:", 8.5, 1.45, 4.5, 0.4, AZUL, 15, bold=True)
add_bullet_box(s,
    ["Cada punto representa una vivienda del conjunto de test.",
     "La línea diagonal = predicción perfecta (predicho = real).",
     "Nuestro modelo se acerca mucho a esa línea diagonal.",
     "Los puntos alejados de la línea son los casos más difíciles de predecir (viviendas muy caras).",
     "El clustering en la zona de $500k se debe al techo artificial del dataset."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=RGBColor(0xE8, 0xF0, 0xFE), text_color=GRIS_TEXTO,
    bullet_char="◉", size_pt=12)

add_numero_slide(s, 17)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — Feature Importance
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("¿Qué variables importan más al modelo?")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "¿Qué variables importan más al modelo? — Feature Importance",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 22, bold=True)

r = add_image_safe(s, img("05_evaluacion_final_img01.png"), 0.4, 1.42, 7.8, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "Insight de negocio:", 8.5, 1.45, 4.5, 0.4, AZUL, 15, bold=True)
add_bullet_box(s,
    ["MedInc (ingresos del vecindario) es el factor más determinante del precio. Explica más del 40% de la importancia total.",
     "Latitude y Longitude aportan conjuntamente ~30% → la ubicación geográfica es crucial.",
     "AveRooms y HouseAge son factores secundarios pero relevantes.",
     "Population y AveBedrms tienen impacto menor en el modelo final.",
     "Para el negocio: si no sabes los ingresos del vecindario, la predicción se degrada significativamente."],
    8.5, 1.9, 4.6, 4.8,
    bg_color=NARANJA_CLARO, text_color=GRIS_TEXTO,
    bullet_char="▶", size_pt=12)

add_numero_slide(s, 18)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19 — Evaluación img02 adicional
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Análisis completo de evaluación")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Análisis completo de evaluación — distribución de errores",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 22, bold=True)

r = add_image_safe(s, img("05_evaluacion_final_img02.png"), 1.5, 1.42, 10.0, 5.3)
if r: img_inserted += 1
else: img_missing += 1

add_text(s, "La distribución de errores nos dice cuánto nos equivocamos y en qué dirección. Un error centrado en cero sin sesgo es ideal.",
         1.5, 6.82, 10.0, 0.55,
         RGBColor(0x88, 0x88, 0x88), 11, italic=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 19)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 20 — Despliegue API (fondo azul)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Etapa 6: El modelo en producción — API REST con FastAPI")
slide_fondo(s, AZUL)
add_rect(s, 0.8, 0.6, 11.73, 0.06, NARANJA)

add_text(s, "Del modelo al servicio: API REST con FastAPI",
         0.8, 0.7, 11.0, 0.75,
         BLANCO, 30, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Etapa 6: El modelo en producción",
         0.8, 1.45, 11.0, 0.45,
         NARANJA, 17, align=PP_ALIGN.CENTER)

# Diagrama de pipeline
steps = ["Notebook\n.ipynb", "Modelo\n.pkl", "FastAPI\napp.py", "Docker\ncontainer", "http://\nlocalhost:8000"]
sx = 0.6
for i, step in enumerate(steps):
    cx = sx + i * 2.55
    step_color = NARANJA if i == 4 else BLANCO
    tc = AZUL if i != 4 else BLANCO
    add_rect(s, cx, 1.95, 2.3, 1.1, step_color)
    add_text(s, step, cx + 0.1, 2.05, 2.1, 0.9,
             tc, 13, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, "→", cx + 2.3, 2.35, 0.25, 0.4,
                 NARANJA, 20, bold=True, align=PP_ALIGN.CENTER)

# JSON snippet
add_rect(s, 0.6, 3.25, 5.5, 2.3, RGBColor(0x1E, 0x1E, 0x2E))
add_text(s, "Respuesta de la API:", 0.75, 3.32, 5.2, 0.35, NARANJA, 12, bold=True)
add_text_multiline(s,
    ['{',
     '  "precio_estimado_usd": 262595,',
     '  "unidad": "USD",',
     '  "modelo_version": "gradient_boosting_v2",',
     '  "confianza": "alta"',
     '}'],
    0.75, 3.65, 5.2, 1.8,
    RGBColor(0x7F, 0xFF, 0x7F), 12)

add_text(s, "Cualquier sistema puede consultar el modelo vía HTTP en milisegundos.\nEl modelo escala a miles de predicciones por segundo.",
         6.4, 3.35, 6.6, 1.2,
         BLANCO, 14, align=PP_ALIGN.CENTER)

# Endpoints
add_rect(s, 6.4, 4.65, 6.6, 1.7, RGBColor(0x00, 0x2A, 0x55))
add_text_multiline(s,
    ["POST  /predict          → predicción individual",
     "POST  /predict/batch   → lote de predicciones",
     "GET   /health           → estado del servicio",
     "GET   /docs             → documentación interactiva"],
    6.55, 4.75, 6.3, 1.5,
    NARANJA, 12)

add_numero_slide(s, 20)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 21 — Interfaz web para usuarios de negocio
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Interfaz web para usuarios de negocio")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Interfaz web para usuarios de negocio",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

add_text(s, "Sin código, sin tecnicismos — http://localhost:8000/ui",
         0.4, 1.25, 12.0, 0.45,
         NARANJA, 16, bold=True)

# Mockup del formulario (texto)
add_rect(s, 0.5, 1.8, 6.0, 5.1, GRIS, border_color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s, "Estimador de Precio de Vivienda",
         0.6, 1.88, 5.8, 0.45,
         AZUL, 14, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, 0.5, 2.3, 6.0, 0.04, NARANJA)

campos = [
    ("Ingresos del vecindario (x $10k)", "3.5"),
    ("Antigüedad de la vivienda (años)",  "15"),
    ("Habitaciones por hogar",            "5"),
    ("Dormitorios por hogar",             "2"),
    ("Población del bloque",              "1200"),
    ("Personas por hogar",                "3"),
    ("Latitud",                           "37.88"),
    ("Longitud",                          "-122.23"),
]
for i, (campo, val) in enumerate(campos):
    cy = 2.4 + i * 0.45
    add_text(s, campo, 0.65, cy, 3.6, 0.38, GRIS_TEXTO, 11)
    add_rect(s, 4.3, cy + 0.03, 1.9, 0.32, BLANCO, RGBColor(0xBB, 0xBB, 0xBB))
    add_text(s, val,   4.4, cy + 0.05, 1.7, 0.28, AZUL, 11)

# Botón
add_rect(s, 1.5, 6.1, 3.5, 0.55, NARANJA)
add_text(s, "Estimar precio", 1.5, 6.16, 3.5, 0.42,
         BLANCO, 14, bold=True, align=PP_ALIGN.CENTER)

# Resultado
add_rect(s, 6.8, 1.8, 6.2, 5.1, RGBColor(0xE8, 0xF8, 0xEE), border_color=VERDE)
add_text(s, "Resultado:", 7.0, 1.95, 5.8, 0.4, GRIS_TEXTO, 13)
add_text(s, "$262,595 USD", 7.0, 2.4, 5.8, 0.8, VERDE, 32, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, 7.5, 3.25, 4.8, 0.04, RGBColor(0xCC, 0xCC, 0xCC))
add_text_multiline(s,
    ["Rango estimado: $240k – $285k",
     "Confianza: Alta",
     "Modelo: GradientBoosting v2",
     "Tiempo de respuesta: 8ms"],
    7.0, 3.35, 5.8, 1.4,
    GRIS_TEXTO, 13)

add_rect(s, 6.8, 5.0, 6.2, 1.75, AZUL)
add_text(s, "Accesible para cualquier usuario de negocio\nsin conocimientos técnicos",
         7.0, 5.2, 5.8, 0.8,
         BLANCO, 14, align=PP_ALIGN.CENTER)
add_text(s, "Disponible 24/7 via navegador web",
         7.0, 6.1, 5.8, 0.45,
         NARANJA, 13, bold=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 21)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 22 — El ciclo MLOps completo (fondo azul)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("El ciclo MLOps completo recorrido")
slide_fondo(s, AZUL)
add_rect(s, 0.8, 0.6, 11.73, 0.06, NARANJA)

add_text(s, "El ciclo MLOps completo recorrido",
         0.8, 0.7, 11.0, 0.75,
         BLANCO, 30, bold=True, align=PP_ALIGN.CENTER)

# 6 pasos en círculo (dispuestos en dos filas de 3)
pasos = [
    ("1", "EDA",           "Entender los datos\ny sus patrones"),
    ("2", "Features",      "Transformar variables\npara el modelo"),
    ("3", "Entrenamiento", "Probar 3 algoritmos\ncon MLflow"),
    ("4", "Selección",     "Elegir el mejor\ncon CV 5-fold"),
    ("5", "Evaluación",    "Test set reservado\nGate de calidad"),
    ("6", "Producción",    "API REST + Docker\nAcceso web"),
]

for i, (num, titulo, desc) in enumerate(pasos):
    col = i % 3
    row = i // 3
    cx = 0.7 + col * 4.15
    cy = 1.55 + row * 2.55
    bg = NARANJA if i % 2 == 0 else BLANCO
    tc = BLANCO if i % 2 == 0 else AZUL
    dc = BLANCO if i % 2 == 0 else GRIS_TEXTO
    add_rect(s, cx, cy, 3.8, 2.1, bg)
    # Número grande
    add_text(s, num, cx + 0.1, cy + 0.1, 0.6, 0.6,
             tc, 28, bold=True)
    add_text(s, titulo, cx + 0.75, cy + 0.12, 2.9, 0.5,
             tc, 18, bold=True)
    add_text_multiline(s, desc.split("\n"),
                       cx + 0.15, cy + 0.68, 3.5, 1.2,
                       dc, 12)
    # Flecha entre pasos
    if i < 5:
        if col < 2:
            add_text(s, "→", cx + 3.8, cy + 0.8, 0.35, 0.5,
                     NARANJA, 20, bold=True)
        else:
            add_text(s, "↓", cx + 1.7, cy + 2.1, 0.4, 0.5,
                     NARANJA, 20, bold=True)

add_text(s, "El monitoreo cierra el ciclo y dispara el re-entrenamiento si el modelo degrada",
         0.7, 7.0, 11.93, 0.4,
         NARANJA, 13, bold=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 22)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 23 — Reproducibilidad total — esto es MLOps
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Reproducibilidad total — esto es MLOps")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Reproducibilidad total — esto es MLOps",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

add_text(s, "Los 4 pilares de este proyecto",
         0.4, 1.25, 12.0, 0.45,
         NARANJA, 16, bold=True)

pilares = [
    ("📦", "Configuración centralizada",
     "config.yaml",
     ["Un único archivo de configuración controla todo el pipeline.",
      "Cambiar un parámetro se propaga automáticamente a todos los notebooks.",
      "Elimina los 'números mágicos' dispersos en el código."]),
    ("🔬", "Tracking de experimentos",
     "MLflow",
     ["Cada entrenamiento queda registrado automáticamente.",
      "Podemos comparar cientos de experimentos visualmente.",
      "El modelo ganador queda versionado y reproducible."]),
    ("🐳", "Containerización",
     "Docker",
     ["El modelo + dependencias empaquetados en un contenedor.",
      "Funciona igual en desarrollo, staging y producción.",
      "Una sola línea: docker-compose up"]),
    ("📓", "Documentación ejecutable",
     "6 Notebooks",
     ["El código es la documentación: comentarios + visualizaciones.",
      "Cualquier persona puede seguir el proceso paso a paso.",
      "Reproducible desde cero en cualquier equipo."]),
]

for i, (icon, titulo, subtitulo, bullets) in enumerate(pilares):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * 6.45
    cy = 1.8 + row * 2.6
    add_rect(s, cx, cy, 6.2, 2.4, GRIS, border_color=AZUL_CLARO)
    add_rect(s, cx, cy, 6.2, 0.55, AZUL)
    add_text(s, f"{icon}  {titulo}", cx + 0.15, cy + 0.08, 4.0, 0.42,
             BLANCO, 15, bold=True)
    add_text(s, subtitulo, cx + 4.2, cy + 0.1, 1.85, 0.38,
             NARANJA, 14, bold=True, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        add_text(s, f"• {b}", cx + 0.15, cy + 0.65 + j * 0.55, 5.9, 0.5,
                 GRIS_TEXTO, 11)

add_rect(s, 0.4, 7.05, 12.53, 0.38, AZUL)
add_text(s, "Cualquier persona puede reproducir este ejemplo en su equipo en menos de 10 minutos",
         0.5, 7.1, 12.3, 0.3,
         NARANJA, 12, bold=True, align=PP_ALIGN.CENTER)

add_numero_slide(s, 23)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 24 — Resumen técnico rápido
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Stack tecnológico del proyecto")
slide_fondo(s, BLANCO)
slide_barra_top(s)

add_text(s, "Stack tecnológico del proyecto",
         0.4, 0.65, 12.0, 0.65,
         AZUL, 26, bold=True)

techs = [
    ("Python 3.11",        "Lenguaje base",              AZUL),
    ("scikit-learn",       "Modelos ML",                 RGBColor(0xF0, 0x8A, 0x00)),
    ("MLflow",             "Experiment Tracking",        RGBColor(0x00, 0x94, 0xF2)),
    ("FastAPI",            "API REST",                   RGBColor(0x05, 0x9A, 0x6B)),
    ("Docker",             "Containerización",           RGBColor(0x00, 0x91, 0xE2)),
    ("Pandas / NumPy",     "Manipulación de datos",      RGBColor(0x15, 0x0D, 0x58)),
    ("Matplotlib / Seaborn","Visualización",             RGBColor(0x11, 0x55, 0x7C)),
    ("YAML config",        "Configuración central",      NARANJA),
    ("Jupyter Notebooks",  "Documentación ejecutable",   RGBColor(0xF3, 0x7B, 0x26)),
]

for i, (tech, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    cx = 0.4 + col * 4.3
    cy = 1.55 + row * 1.6
    add_rect(s, cx, cy, 4.1, 1.45, GRIS, border_color=color)
    add_rect(s, cx, cy, 4.1, 0.45, color)
    add_text(s, tech, cx + 0.12, cy + 0.06, 3.86, 0.36,
             BLANCO, 14, bold=True)
    add_text(s, desc, cx + 0.12, cy + 0.55, 3.86, 0.75,
             GRIS_TEXTO, 12)

add_numero_slide(s, 24)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 25 — CIERRE
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("CIERRE — Gracias / Preguntas")
slide_fondo(s, AZUL)
add_rect(s, 0.8, 0.6, 11.73, 0.06, NARANJA)

add_text(s, "¡Gracias!",
         0.8, 0.75, 7.5, 1.0,
         BLANCO, 54, bold=True)

add_text(s, "Preguntas y debate",
         0.8, 1.75, 7.5, 0.6,
         NARANJA, 28)

add_rect(s, 0.8, 2.55, 7.5, 0.05, BLANCO)

add_text_multiline(s,
    ["Prueba el modelo en vivo:",
     "http://localhost:8000/ui",
     "",
     "Documentación interactiva:",
     "http://localhost:8000/docs",
     "",
     "Repositorio del proyecto:",
     "github.com/tu-org/mlops-ciclo-vida"],
    0.8, 2.7, 7.0, 3.8,
    BLANCO, 15)

# Foto instructor
if os.path.exists(str(FOTO_INS)):
    s.shapes.add_picture(str(FOTO_INS), Inches(9.5), Inches(1.5), Inches(3.0), Inches(3.0))
    img_inserted += 1
    add_text(s, "Instructor", 9.5, 4.55, 3.0, 0.4,
             BLANCO, 13, bold=True, align=PP_ALIGN.CENTER)

add_text_multiline(s,
    ["Dataset: California Housing (sklearn)",
     "Modelo: GradientBoosting — RMSE 0.43 / R² 0.84",
     "6 notebooks + API + Docker = MLOps completo"],
    0.8, 6.65, 8.5, 0.75,
    RGBColor(0xAA, 0xCC, 0xFF), 11)

add_numero_slide(s, 25)


# ─────────────────────────────────────────────────────────────────────────────
# GUARDAR
# ─────────────────────────────────────────────────────────────────────────────
OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))

# ─────────────────────────────────────────────────────────────────────────────
# REPORTE
# ─────────────────────────────────────────────────────────────────────────────
size_kb = OUTPUT.stat().st_size / 1024
print(f"\n{'='*60}")
print(f"PRESENTACION GENERADA: {OUTPUT}")
print(f"Tamano: {size_kb:.1f} KB  ({OUTPUT.stat().st_size:,} bytes)")
print(f"{'='*60}")
print(f"Imagenes de notebooks insertadas : {img_inserted}")
print(f"Imagenes no encontradas          : {img_missing}")
print(f"Total slides                     : {len(slide_titles)}")
print(f"{'='*60}")
print("\nSlides generadas:")
for i, t in enumerate(slide_titles, 1):
    print(f"  {i:2d}. {t}")
